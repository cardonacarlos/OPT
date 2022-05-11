import tkinter as tk
import webbrowser
import requests
import pandas as pd
import numpy as np
import json
import math
import re
# import pptx
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN
from logic import *


#header for accessing API
headers = {
    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/95.0.4638.54 Safari/537.36',
    'From': 'ccardona2000@gmail.com' 
}

# welcome user
def write_slogan():
   print('Welcome to the application!')

# find a company's SEC report from ticker
def hyperlink(CIK):
    url = 'https://www.sec.gov/edgar/browse/?CIK=' + CIK
    webbrowser.open_new_tab(url)

# read ticker to cik file
# refer to link https://www.sec.gov/files/company_tickers.json
def tick_to_cik(ticker):
    url = 'https://www.sec.gov/files/company_tickers.json'
    response = requests.get(url, headers=headers)
    if not response.ok:
        print('Retrying...')
    else:
        data = response.json()
        for tick in data.values():
            if tick['ticker'].lower() == ticker.lower():
                print(tick['title'])
                return str(tick['cik_str']).zfill(10)
            
#list of items that we want to observe
items = """
Assets
CommonStockSharesAuthorized
CommonStockSharesIssued
CommonStockSharesOutstanding
CommonStockDividendsPerShareDeclared
CostOfGoodsAndServicesSold
Dividends
EarningsPerShareDiluted
Liabilities
LiabilitiesAndStockholdersEquity
LongTermDebt
LongTermDebtCurrent
LongTermDebtMaturitiesRepaymentsOfPrincipalAfterYearFive
OperatingExpenses
OperatingIncomeLoss
PreferredStockSharesAuthorized
Revenues
StockholdersEquity
WeightedAverageNumberOfSharesOutstandingBasic
NetIncomeLoss
"""
items = items.split()
global info
info = ""




# access SEC Edgar API and return facts
def company_facts(string, items, info, info_label):
    
    if string:
        info = string.upper() + " Company Facts\n"
        url = "https://data.sec.gov/api/xbrl/companyfacts/CIK"  + tick_to_cik(string) + ".json"
    
        response = requests.get(url, headers=headers)
        
        #ensure that connect works
        if not response.ok:
            print('Retry')
        else:
            #load json data
            response = response.text
            data = json.loads(response)

            #dictionary for wanted metric
            global metric_dict
            metric_dict = {}


            for key, val in data['facts']['us-gaap'].items():
                unit_df = pd.DataFrame()

                for unit, unit_val in val['units'].items():
                    val_data = val['units'][unit]
                    val_data_df = pd.DataFrame(unit_val)
                    val_data_df['unit'] = unit
                    
                    
                    # deprecated warning
                    # unit_df = unit_df.append(val_data_df)
                    unit_df = pd.concat([unit_df, val_data_df])
                    metric_dict[key] = unit_df
            
            print([key for key in metric_dict.keys()])
            
            
            for item in items:
                
                if item in metric_dict:
                    metric = metric_dict[item]
                    ten_k = metric[metric['form'] == '10-K']
                    date = max(ten_k['filed'])
                    ten_k = ten_k[(ten_k['filed'] == date) & (ten_k['form'] == '10-K') & (ten_k['frame'].isna())]

                    #if the value for the requested metric is not null, then print
                    #else find the latest metric for it available
                    if ten_k['val'].empty:
                        ten_k = metric
                        date = max(ten_k['filed'])
                        ten_k = ten_k[(ten_k['filed']==date)]   
                        print(item + ": Value: ")
                        print(ten_k.iloc[0]['val'])
                        info = info + " ".join(re.sub(r"([A-Z])", r" \1", str(item)).split()) + ": " + "{:,}".format(ten_k.iloc[0]['val']) + "\n"
                        info_label.config(text=info)

                    else:
                        print(item + ": Value: ")
                        print(ten_k.iloc[0]['val'])
                        info = info + " ".join(re.sub(r"([A-Z])", r" \1", str(item)).split()) + ": " + "{:,}".format(ten_k.iloc[0]['val']) +"\n"
                        info_label.config(text=info)
        print("\n===================================================")
        print(info)
    else:
        print('No Input')
        
# company_facts("aapl", items, info, info_label=None)
    
#command that makes powerpoint from financial information
def make_powerpoint():

    #split financial info into array
    split_info = info.split("\n")

    #initialize powerpoint
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Company Facts for " + string.upper()

    #make shape for each info
    shapes = slide.shapes
    left = top = height = Inches(0.5)
    width = Inches(0.7)
    for text in range(len(split_info)):

        #shape configuration
        shape = shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
        )

        #shape filling
        fill = shape.fill 
        fill.solid()
        fill.fore_color.rgb = RGBColor(255,255,255)

        #shape outline
        line = shape.line
        line.color.rgb = RGBColor(0,0,0)

        #shape text
        text_frame = shape.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = split_info[text]

        #shape text font
        font = run.font
        font.name = 'Arial Narrow'
        font.size = Pt(6)
        font.color.rgb = RGBColor(0,0,0)
        left = left + width

    pptx_name = string.upper() + ".pptx"
    prs.save(pptx_name)

    #open powerpoint file that was just created

    cwd = os.getcwd()
    path = cwd + "\\" + pptx_name
    os.startfile(path)
    
# create a checklist of unused items
def check():
    if metric_dict:
        print("Metric_Dict Exists")
        print(metric_dict)
        print(metric_dict.keys())
    else:
        print("No Change")